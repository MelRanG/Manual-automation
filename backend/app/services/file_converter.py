import asyncio
import io
from pathlib import Path

from app.config import settings


STATIC_IMAGES_DIR = Path(__file__).resolve().parent.parent.parent / "static" / "images"
STATIC_IMAGES_DIR.mkdir(parents=True, exist_ok=True)


def _use_s3() -> bool:
    return bool(settings.uploads_s3_bucket)


def _upload_image_to_s3(image_bytes: bytes, document_id: str, filename: str) -> str:
    import boto3

    prefix = settings.uploads_s3_prefix.strip("/")
    key = f"{prefix}/images/{document_id}/{filename}" if prefix else f"images/{document_id}/{filename}"
    client = boto3.client("s3", region_name=settings.aws_region)
    client.put_object(
        Bucket=settings.uploads_s3_bucket,
        Key=key,
        Body=image_bytes,
        ServerSideEncryption="AES256",
    )
    return f"https://{settings.uploads_s3_bucket}.s3.{settings.aws_region}.amazonaws.com/{key}"


def _save_image(image_bytes: bytes, document_id: str, filename: str, static_dir: Path) -> str:
    if _use_s3():
        return _upload_image_to_s3(image_bytes, document_id, filename)
    img_dir = static_dir / document_id
    img_dir.mkdir(parents=True, exist_ok=True)
    (img_dir / filename).write_bytes(image_bytes)
    return f"/static/images/{document_id}/{filename}"


async def convert_to_markdown(
    file_bytes: bytes,
    filename: str,
    document_id: str,
    static_dir: Path | None = None,
) -> str:
    return await asyncio.to_thread(
        _convert_sync, file_bytes, filename, document_id, static_dir or STATIC_IMAGES_DIR
    )


def _convert_sync(
    file_bytes: bytes,
    filename: str,
    document_id: str,
    static_dir: Path,
) -> str:
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _pdf_to_markdown(file_bytes, document_id, static_dir)
    if ext == ".pptx":
        return _pptx_to_markdown(file_bytes, document_id, static_dir)
    if ext == ".docx":
        return _docx_to_markdown(file_bytes, document_id, static_dir)
    if ext in (".xlsx", ".xls"):
        return _xlsx_to_markdown(file_bytes)
    return file_bytes.decode("utf-8", errors="replace")


def _pdf_to_markdown(file_bytes: bytes, document_id: str, static_dir: Path) -> str:
    import fitz

    parts: list[str] = []
    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num, page in enumerate(pdf, start=1):
        text = page.get_text().strip()
        if text:
            parts.append(f"## 페이지 {page_num}\n\n{text}")

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf.extract_image(xref)
            img_bytes = base_image["image"]
            ext = base_image["ext"]
            img_filename = f"page{page_num}_img{img_index}.{ext}"
            url = _save_image(img_bytes, document_id, img_filename, static_dir)
            parts.append(f"![이미지]({url})")

    pdf.close()
    return "\n\n".join(parts)


def _pptx_to_markdown(file_bytes: bytes, document_id: str, static_dir: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"## 슬라이드 {slide_num}"]
        img_idx = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_filename = f"slide{slide_num}_img{img_idx}.png"
                url = _save_image(shape.image.blob, document_id, img_filename, static_dir)
                slide_parts.append(f"![이미지]({url})")
                img_idx += 1
        parts.append("\n\n".join(slide_parts))

    return "\n\n---\n\n".join(parts)


def _docx_to_markdown(file_bytes: bytes, document_id: str, static_dir: Path) -> str:
    import mammoth
    import markdownify

    img_counter = 0

    def handle_image(image):
        nonlocal img_counter
        with image.open() as f:
            img_bytes = f.read()
        content_type = image.content_type or "image/png"
        ext = content_type.split("/")[-1]
        img_filename = f"img{img_counter}.{ext}"
        img_counter += 1
        url = _save_image(img_bytes, document_id, img_filename, static_dir)
        return {"src": url}

    result = mammoth.convert_to_html(
        io.BytesIO(file_bytes),
        convert_image=mammoth.images.img_element(handle_image),
    )
    return markdownify.markdownify(result.value, heading_style="ATX")


def _xlsx_to_markdown(file_bytes: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"## {sheet.title}\n")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else "" for c in rows[0]]
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("")
    wb.close()
    return "\n".join(parts)
