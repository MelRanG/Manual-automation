import io
import uuid
from pathlib import Path

import pytest

from app.services.file_converter import convert_to_markdown


@pytest.mark.asyncio
async def test_txt_passthrough(tmp_path: Path):
    result = await convert_to_markdown(
        file_bytes=b"hello world",
        filename="test.txt",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "hello world" in result


@pytest.mark.asyncio
async def test_md_passthrough(tmp_path: Path):
    result = await convert_to_markdown(
        file_bytes=b"# Hello",
        filename="test.md",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "# Hello" in result


@pytest.mark.asyncio
async def test_xlsx_to_markdown_table(tmp_path: Path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["이름", "나이"])
    ws.append(["홍길동", 30])
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    result = await convert_to_markdown(
        file_bytes=buf.read(),
        filename="test.xlsx",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "이름" in result
    assert "홍길동" in result
    assert "|" in result


@pytest.mark.asyncio
async def test_docx_to_markdown(tmp_path: Path):
    from docx import Document as DocxDocument
    doc = DocxDocument()
    doc.add_heading("제목", level=1)
    doc.add_paragraph("본문 내용입니다.")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    result = await convert_to_markdown(
        file_bytes=buf.read(),
        filename="test.docx",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "제목" in result
    assert "본문 내용" in result


@pytest.mark.asyncio
async def test_pdf_to_markdown(tmp_path: Path):
    import fitz
    pdf_doc = fitz.open()
    page = pdf_doc.new_page()
    rect = fitz.Rect(50, 100, 400, 200)
    page.insert_htmlbox(rect, "<p>PDF 테스트 내용</p>")
    buf = io.BytesIO(pdf_doc.tobytes())

    result = await convert_to_markdown(
        file_bytes=buf.read(),
        filename="test.pdf",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "PDF 테스트 내용" in result


@pytest.mark.asyncio
async def test_pptx_to_markdown(tmp_path: Path):
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "슬라이드 제목"
    slide.placeholders[1].text = "슬라이드 내용"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = await convert_to_markdown(
        file_bytes=buf.read(),
        filename="test.pptx",
        document_id=str(uuid.uuid4()),
        static_dir=tmp_path,
    )
    assert "슬라이드 제목" in result


@pytest.mark.asyncio
async def test_image_extracted_to_static(tmp_path: Path):
    """PDF에서 이미지가 추출되면 static_dir에 저장되어야 한다."""
    import fitz
    from PIL import Image as PILImage

    # 1x1 PNG 이미지를 포함하는 PDF 생성
    img = PILImage.new("RGB", (10, 10), color=(255, 0, 0))
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)

    pdf_doc = fitz.open()
    page = pdf_doc.new_page()
    text_rect = fitz.Rect(50, 100, 400, 140)
    page.insert_htmlbox(text_rect, "<p>이미지 포함 PDF</p>")
    rect = fitz.Rect(50, 150, 150, 250)
    page.insert_image(rect, stream=img_buf.read())
    buf = io.BytesIO(pdf_doc.tobytes())

    doc_id = str(uuid.uuid4())
    result = await convert_to_markdown(
        file_bytes=buf.read(),
        filename="test_with_image.pdf",
        document_id=doc_id,
        static_dir=tmp_path,
    )
    # 이미지 파일이 static_dir에 생성되어야 함
    image_dir = tmp_path / doc_id
    assert image_dir.exists()
    assert any(image_dir.iterdir())
    # Markdown에 이미지 참조가 포함되어야 함
    assert "![" in result
    assert f"/static/images/{doc_id}/" in result
